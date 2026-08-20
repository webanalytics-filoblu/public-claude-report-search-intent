#!/usr/bin/env python3
"""
Genera il pptx completo del deck, interamente offline (nessuna chiamata Google): sostituisce
build_slides.py, che lavorava in due fasi su una presentazione remota (duplica slide vuote via
Slides API, poi Claude scriveva il testo con replaceAllText scoped sulla slide nuova).

Qui non esiste piu' un editing incrementale post-caricamento (l'MCP Drive non offre un
batchUpdate su un file esistente): tutto il contenuto testuale/editoriale va deciso PRIMA di
chiamare questo script e passato via --slide-content (vedi SKILL.md, Step 6b-nuovo). Lo script
fa, in un solo passaggio:
  1. Duplica, per ogni cluster/colonna extra, la slide-stencil scelta in --stencil-map /
     --breakdown-stencil-map (qui identificata per indice di slide nel template, non piu' per
     objectId Slides API — python-pptx non ha "duplicate slide" nativo: si clona l'XML della
     slide sorgente).
  2. Sostituisce ogni placeholder di testo "####..." col contenuto di --slide-content.
  3. Applica il colore RGB richiesto per il delta YoY (verde/rosso), se presente in
     --slide-content.
  4. Trova l'immagine piu' grande sulla slide (stessa euristica area-massima di prima) e la
     sostituisce con un grafico nativo pptx (torta) costruito dai valori di --charts-meta, o
     con un'immagine statica se fornita in --brand-images.
  5. Sostituisce il placeholder del nome brand in tutto il deck.

Claude carica poi il pptx risultante su Drive con mcp__claude_ai_Google_Drive__create_file
(content_mime_type pptx, conversione automatica in Google Slides).

Uso:
    python scripts/build_slides_pptx.py \
        --template-pptx runs/<slug>/template/slide_template.pptx \
        --charts-meta runs/<slug>/charts_meta.json \
        --stencil-map runs/<slug>/stencil_map.json \
        --slide-content runs/<slug>/slide_content.json \
        --new-brand-text "Yamamay" \
        --pptx-out runs/<slug>/output/deck.pptx \
        --slides-meta-out runs/<slug>/slides_meta.json \
        [--breakdown-stencil-map runs/<slug>/breakdown_stencil_map.json] \
        [--cluster-overview-slide-index 4] \
        [--brand-images cover=path/to/cover.jpg,main_insights=path/to/main.jpg]
"""

import argparse
import copy
import json
from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.dml.color import RGBColor
from pptx.util import Emu

DELTA_COLORS = {
    "green": RGBColor(0x00, 0x80, 0x00),
    "red": RGBColor(0xC0, 0x00, 0x00),
}


def duplicate_slide(prs, source_index):
    """Clona la slide a source_index e la aggiunge in coda al deck. python-pptx non espone
    un'API di duplicazione: si copia l'XML dell'elemento <p:sld> sorgente e lo si registra
    come nuova slide, stessa tecnica usata per gli stessi scopi in altri progetti python-pptx
    (nessuna libreria terza aggiunta). Le relationship (immagini embedded, layout) vengono
    riusate dalla slide sorgente."""
    source = prs.slides[source_index]
    blank_slide_layout = source.slide_layout
    new_slide = prs.slides.add_slide(blank_slide_layout)

    # Rimuove i placeholder vuoti creati automaticamente dal layout, poi copia gli shape reali
    # della slide sorgente (compresi quelli con testo/immagini) — evita doppioni.
    for shape in list(new_slide.shapes):
        shape._element.getparent().remove(shape._element)

    for shape in source.shapes:
        new_slide.shapes._spTree.append(copy.deepcopy(shape._element))

    if source.has_notes_slide:
        new_slide.notes_slide.notes_text_frame.text = source.notes_slide.notes_text_frame.text

    return new_slide


def replace_placeholder_text(slide, replacements):
    """Sostituisce, per ogni shape con testo, ogni occorrenza esatta di una chiave di
    `replacements` con il suo valore. Non tenta un replace parziale su singoli run (i
    placeholder "####..." sono runs unici nel template, verificato dal vivo prima di questo
    script) — se un placeholder e' diviso su piu' run, va normalizzato a monte nel template."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for placeholder, value in replacements.items():
                    if placeholder in run.text:
                        run.text = run.text.replace(placeholder, value)


def apply_delta_color(slide, delta_markers):
    """delta_markers: {"testo esatto del run": "green"|"red"} — applica il colore al run che
    contiene esattamente quel testo (dopo la sostituzione placeholder, quindi il testo cercato
    e' il valore GIA' sostituito, non il placeholder originale)."""
    if not delta_markers:
        return
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs:
                for marker_text, color_name in delta_markers.items():
                    if marker_text in run.text and color_name in DELTA_COLORS:
                        run.font.color.rgb = DELTA_COLORS[color_name]


def largest_picture(slide):
    pictures = [sh for sh in slide.shapes if sh.shape_type == 13]
    if not pictures:
        return None
    return max(pictures, key=lambda sh: sh.width * sh.height)


def replace_with_pie_chart(slide, placeholder_shape, categories, values, chart_title):
    left, top, width, height = (
        placeholder_shape.left, placeholder_shape.top,
        placeholder_shape.width, placeholder_shape.height,
    )
    placeholder_shape._element.getparent().remove(placeholder_shape._element)

    chart_data = CategoryChartData()
    chart_data.categories = categories
    chart_data.add_series(chart_title, values)

    graphic_frame = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE, Emu(left), Emu(top), Emu(width), Emu(height), chart_data,
    )
    graphic_frame.chart.has_title = True
    graphic_frame.chart.chart_title.text_frame.text = chart_title


def replace_with_image(slide, placeholder_shape, image_path):
    left, top, width, height = (
        placeholder_shape.left, placeholder_shape.top,
        placeholder_shape.width, placeholder_shape.height,
    )
    placeholder_shape._element.getparent().remove(placeholder_shape._element)
    slide.shapes.add_picture(image_path, Emu(left), Emu(top), Emu(width), Emu(height))


def attach_chart_or_image(slide, records, value_label, chart_title, brand_image_path, label):
    placeholder_shape = largest_picture(slide)
    if not placeholder_shape:
        print(f"ATTENZIONE: nessuna immagine trovata sulla slide di '{label}' — "
              f"nessun grafico/immagine inserito, va aggiunto a mano.")
        return

    if brand_image_path:
        replace_with_image(slide, placeholder_shape, brand_image_path)
        return

    if not records:
        print(f"ATTENZIONE: nessun dato per il grafico di '{label}' — placeholder lasciato invariato.")
        return

    categories = [r[value_label] for r in records]
    values = [r["Volume Totale"] for r in records]
    replace_with_pie_chart(slide, placeholder_shape, categories, values, chart_title)


def build_slide(prs, stencil_index, content, chart_records, value_label, chart_title, brand_image_path, label):
    slide = duplicate_slide(prs, stencil_index)

    replacements = dict(content.get("replacements", {}))
    replace_placeholder_text(slide, replacements)
    apply_delta_color(slide, content.get("delta_colors", {}))
    attach_chart_or_image(slide, chart_records, value_label, chart_title, brand_image_path, label)

    return slide


def main():
    parser = argparse.ArgumentParser(description="Genera il pptx completo del deck (offline, nessuna chiamata Google)")
    parser.add_argument("--template-pptx", required=True)
    parser.add_argument("--charts-meta", required=True)
    parser.add_argument("--stencil-map", required=True,
                         help='JSON {"label cluster": <indice slide-stencil nel template>}')
    parser.add_argument("--breakdown-stencil-map",
                         help='JSON {"colonna extra": <indice slide-stencil>} (facoltativo)')
    parser.add_argument("--slide-content", required=True,
                         help='JSON {"clusters": {label: {"replacements": {...}, "delta_colors": {...}}}, '
                              '"breakdowns": {...}, "fixed": {...}} — testo editoriale scritto da Claude '
                              'PRIMA di generare il pptx (vedi SKILL.md Step 6b-nuovo)')
    parser.add_argument("--brand-placeholder-text", default="Nome brand in MAIUSCOLO")
    parser.add_argument("--new-brand-text", required=True)
    parser.add_argument("--cluster-overview-slide-index", type=int,
                         help="Indice (0-based) della slide fissa 'cluster overview' nel template, se presente")
    parser.add_argument("--brand-images",
                         help='"chiave=path,chiave2=path2" — path locali da inserire al posto del '
                              'placeholder immagine sulle slide FISSE indicate (cover, main_insights, ...)')
    parser.add_argument("--pptx-out", required=True)
    parser.add_argument("--slides-meta-out", required=True)
    args = parser.parse_args()

    charts_meta = json.loads(Path(args.charts_meta).read_text(encoding="utf-8"))
    stencil_map = json.loads(Path(args.stencil_map).read_text(encoding="utf-8"))
    breakdown_stencil_map = (
        json.loads(Path(args.breakdown_stencil_map).read_text(encoding="utf-8"))
        if args.breakdown_stencil_map else {}
    )
    slide_content = json.loads(Path(args.slide_content).read_text(encoding="utf-8"))
    brand_images = dict(
        pair.split("=", 1) for pair in args.brand_images.split(",")
    ) if args.brand_images else {}

    unknown_stencils = [c for c in stencil_map if c not in charts_meta["clusters"]]
    if unknown_stencils:
        raise SystemExit(f"--stencil-map cita cluster non presenti in charts_meta.json: {', '.join(unknown_stencils)}")
    unknown_breakdown_stencils = [c for c in breakdown_stencil_map if c not in charts_meta.get("breakdowns", {})]
    if unknown_breakdown_stencils:
        raise SystemExit(f"--breakdown-stencil-map cita colonne non presenti in charts_meta.json: "
                          f"{', '.join(unknown_breakdown_stencils)}")

    prs = Presentation(args.template_pptx)

    if args.cluster_overview_slide_index is not None:
        overview = charts_meta.get("overview")
        overview_content = slide_content.get("fixed", {}).get("cluster_overview", {})
        overview_slide = prs.slides[args.cluster_overview_slide_index]
        replace_placeholder_text(overview_slide, overview_content.get("replacements", {}))
        apply_delta_color(overview_slide, overview_content.get("delta_colors", {}))
        if overview:
            attach_chart_or_image(
                overview_slide, overview.get("subClusters", []), "Cluster",
                "Volume di ricerca per Cluster", brand_images.get("cluster_overview"), "cluster overview",
            )

    for name, content in slide_content.get("fixed", {}).items():
        if name == "cluster_overview":
            continue
        # Le slide fisse (cover, brand-secco, main insights, warnings) esistono gia' nel
        # template: qui si applica solo testo/colore/immagine, senza duplicare nulla.
        idx = content.get("slide_index")
        if idx is None:
            continue
        fixed_slide = prs.slides[idx]
        replace_placeholder_text(fixed_slide, content.get("replacements", {}))
        apply_delta_color(fixed_slide, content.get("delta_colors", {}))
        if name in brand_images:
            placeholder_shape = largest_picture(fixed_slide)
            if placeholder_shape:
                replace_with_image(fixed_slide, placeholder_shape, brand_images[name])

    slides_meta = {"clusters": {}, "breakdowns": {}}

    for label, stencil_index in stencil_map.items():
        info = charts_meta["clusters"][label]
        content = slide_content.get("clusters", {}).get(label, {})
        slide = build_slide(
            prs, stencil_index, content, info.get("subClusters", []), "Sotto Cluster",
            f"{label} — Ripartizione Sotto Cluster", None, f"cluster '{label}'",
        )
        replace_placeholder_text(slide, {args.brand_placeholder_text: args.new_brand_text})
        slides_meta["clusters"][label] = {"slideIndex": len(prs.slides) - 1}
        print(f"Slide creata per cluster '{label}' (indice {len(prs.slides) - 1})")

    for column, stencil_index in breakdown_stencil_map.items():
        info = charts_meta["breakdowns"][column]
        content = slide_content.get("breakdowns", {}).get(column, {})
        slide = build_slide(
            prs, stencil_index, content, info.get("values", []), column,
            f"Distribuzione per {column}", None, f"breakdown '{column}'",
        )
        replace_placeholder_text(slide, {args.brand_placeholder_text: args.new_brand_text})
        slides_meta["breakdowns"][column] = {"slideIndex": len(prs.slides) - 1}
        print(f"Slide creata per breakdown '{column}' (indice {len(prs.slides) - 1})")

    # Sostituzione finale del placeholder brand su TUTTO il deck (incluse le slide fisse
    # gia' processate sopra, che potrebbero non averlo ricevuto se non presente in slide_content).
    for slide in prs.slides:
        replace_placeholder_text(slide, {args.brand_placeholder_text: args.new_brand_text})

    pptx_out = Path(args.pptx_out)
    pptx_out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(pptx_out)

    Path(args.slides_meta_out).write_text(json.dumps(slides_meta, indent=2, ensure_ascii=False), encoding="utf-8")
    print(f"\npptx: {pptx_out}")
    print(f"slides_meta: {args.slides_meta_out}")


if __name__ == "__main__":
    main()
