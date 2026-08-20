#!/usr/bin/env python3
"""
Diagnostica READ-ONLY del template Slides, offline: sostituisce la parte "Slides" di
inspect_template.py, che leggeva il template remoto via Slides API.

Il template va prima scaricato/esportato in locale con il tool MCP
mcp__claude_ai_Google_Drive__download_file_content (fileId=GOOGLE_SLIDE_TEMPLATE_ID,
exportMimeType=application/vnd.openxmlformats-officedocument.presentationml.presentation), da
Claude nella conversazione — questo script non fa alcuna chiamata di rete, solo lettura locale.

Come prima, il template contiene testi placeholder marcati col prefisso "####" (corpo slide e
note relatore) — rileggilo SEMPRE fresco prima di ogni run (l'export puo' cambiare se qualcuno
modifica il template su Google Slides nel frattempo): non fidarti di un'ispezione precedente.

Stampa, per ogni slide: indice (usato al posto dell'objectId Slides API — stabile finche' il
template non viene riordinato; per un riferimento piu' stabile, preferisci nello stencil_map il
testo placeholder univoco letto qui, non l'indice), note relatore, testo di ogni shape con
placeholder, ed evidenzia l'immagine con l'area maggiore (il grafico-segnaposto, per la stessa
euristica di area-massima usata prima per distinguere logo/grafico).

Non modifica nulla: solo lettura e stampa.

Uso:
    python scripts/inspect_template_pptx.py --pptx runs/<slug>/template/slide_template.pptx
"""

import argparse

from pptx import Presentation


def shape_text(shape) -> str:
    if not shape.has_text_frame:
        return ""
    return shape.text_frame.text.strip()


def largest_picture(slide):
    pictures = [sh for sh in slide.shapes if sh.shape_type == 13]  # MSO_SHAPE_TYPE.PICTURE
    if not pictures:
        return None
    return max(pictures, key=lambda sh: sh.width * sh.height)


def main():
    parser = argparse.ArgumentParser(description="Ispeziona (sola lettura) un template Slides esportato in pptx")
    parser.add_argument("--pptx", required=True, help="Percorso locale del pptx esportato dal template Google Slides")
    args = parser.parse_args()

    prs = Presentation(args.pptx)
    print(f"File: {args.pptx}")
    print(f"Slide totali: {len(prs.slides)}")

    for i, slide in enumerate(prs.slides):
        print(f"\n--- Slide {i} ---")
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                print(f"  Note relatore: {notes!r}")

        for sh in slide.shapes:
            if sh.shape_type == 13:
                continue
            text = shape_text(sh)
            if text:
                print(f"  [shape id={sh.shape_id}] text={text!r}")

        biggest = largest_picture(slide)
        if biggest:
            print(f"  [immagine piu' grande] shape id={biggest.shape_id} "
                  f"size={biggest.width}x{biggest.height} pos=({biggest.left},{biggest.top})")


if __name__ == "__main__":
    main()
